from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {"title": "Using Skills - Overview", "bullets": [
        "What is a skill? - small reusable function/module",
        "When to use skills: composition, testing, reuse",
        "Goals of this tutorial: build, call, test, extend"
    ]},
    {"title": "Skill Example: greet", "bullets": [
        "Simple function: greet(name) -> str",
        "Shows basic input/output and testing",
        "Exercise: modify to support languages"
    ]},
    {"title": "Skill Example: calculator", "bullets": [
        "add(a,b), divide(a,b)",
        "Edge cases: divide by zero",
        "Exercise: add multiply and subtract"
    ]},
    {"title": "I/O Skill: fetch_url", "bullets": [
        "Performs network request (requests)",
        "Handle timeouts and errors",
        "Exercise: add JSON parsing and caching"
    ]},
    {"title": "File Skill: summarize_file", "bullets": [
        "Reads a file and returns stats",
        "Permissions and encoding considerations",
        "Exercise: support binary files or large files"
    ]},
    {"title": "Organizing Skills", "bullets": [
        "Group related skills in a module (skills.py)",
        "Provide a dispatcher or registry",
        "Keep pure logic separate from I/O"
    ]},
    {"title": "Dispatcher Pattern", "bullets": [
        "SKILLS = {name: fn}",
        "call_skill(name, *args, **kwargs)",
        "Benefits: dynamic calling, plugin-friendly"
    ]},
    {"title": "Testing Skills", "bullets": [
        "Write unit tests for pure functions",
        "Mock network and file I/O",
        "CI integration: run tests on push"
    ]},
    {"title": "Error Handling", "bullets": [
        "Raise clear exceptions (ValueError)",
        "Return structured errors for callers",
        "Document failure modes"
    ]},
    {"title": "Logging and Observability", "bullets": [
        "Use logging module, not prints",
        "Add correlation ids for traces",
        "Expose metrics where useful"
    ]},
    {"title": "Stateful vs Stateless Skills", "bullets": [
        "Prefer stateless functions for testability",
        "Stateful skills can hold caches or sessions",
        "Persist state carefully and explicitly"
    ]},
    {"title": "Async Skills", "bullets": [
        "Use async/await for concurrency",
        "Provide both sync and async wrappers if needed",
        "Beware of blocking I/O in async code"
    ]},
    {"title": "Dependency Management", "bullets": [
        "Pin versions in requirements.txt",
        "Document runtime dependencies",
        "Use virtual environments"
    ]},
    {"title": "Security Considerations", "bullets": [
        "Sanitize inputs for I/O skills",
        "Avoid executing untrusted code",
        "Handle secrets via env vars or vaults"
    ]},
    {"title": "Packaging Skills", "bullets": [
        "Make reusable packages (pip) if needed",
        "Provide clear APIs and docs",
        "Semantic versioning for compatibility"
    ]},
    {"title": "Documentation", "bullets": [
        "Docstrings for each skill",
        "README with examples",
        "API examples and quickstart"
    ]},
    {"title": "Examples: greet Usage", "bullets": [
        "Python: greet('Alice')",
        "CLI wrapper: skill-run greet Alice",
        "HTTP wrapper: /skills/greet?name=Bob"
    ]},
    {"title": "Examples: calculator Usage", "bullets": [
        "Direct call in code or REPL",
        "Validate numeric inputs",
        "Use for pipelines and data transforms"
    ]},
    {"title": "Examples: fetch_url Usage", "bullets": [
        "Return status & snippet",
        "Handle JSON vs HTML",
        "Respect robots and rate limits"
    ]},
    {"title": "Advanced: Composing Skills", "bullets": [
        "Chain skills: fetch -> parse -> summarize",
        "Use small focused skills for composition",
        "Maintain clear contracts between steps"
    ]},
    {"title": "Advanced: Plugins", "bullets": [
        "Load skills dynamically from folders",
        "Register via entry points or config",
        "Hot-reload for development"
    ]},
    {"title": "Debugging Skills", "bullets": [
        "Add detailed logs in dev",
        "Reproduce with unit tests",
        "Use small harness scripts to exercise skills"
    ]},
    {"title": "Performance", "bullets": [
        "Benchmark hot paths",
        "Cache results where safe",
        "Avoid expensive I/O in tight loops"
    ]},
    {"title": "Versioning and Compatibility", "bullets": [
        "Document breaking changes",
        "Support deprecation policies",
        "Provide migration guides"
    ]},
    {"title": "Deployment Patterns", "bullets": [
        "Bundle skills with app or run as separate service",
        "Use containers for isolation",
        "Monitor health and errors"
    ]},
    {"title": "Exercises", "bullets": [
        "Implement a new skill: translate(text, lang)",
        "Add tests and CI for it",
        "Document usage and edge cases"
    ]},
    {"title": "Resources", "bullets": [
        "python-pptx docs (used to generate this)",
        "requests, pytest, logging docs",
        "Design patterns for small services"
    ]},
    {"title": "Q & A", "bullets": [
        "Questions to try: when to refactor a skill?",
        "How to test I/O-heavy skills?",
        "Next steps for learning"
    ]},
    {"title": "Summary", "bullets": [
        "Prefer small, testable skills",
        "Document and observe them",
        "Practice by building and composing"
    ]}
]

prs = Presentation()
# Title slide layout index may vary; use layout 0 for title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Skills: How to design and use them"
subtitle = slide.placeholders[1]
subtitle.text = "Auto-generated tutorial slides"

for s in slides:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = s['title']
    # content placeholder may be the second shape
    body = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title:
            body = shape
            break
    if body is None:
        # create textbox fallback
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        body = txBox

    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(s['bullets']):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

out_path = "skills_presentation.pptx"
prs.save(out_path)
print(f"Saved presentation to {out_path}")
